from flask import Flask, request, jsonify, send_file, render_template
from pptx import Presentation
from pptx.util import Inches
import requests
import time
import os
import uuid

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

API_TOKEN = os.getenv("API_TOKEN")

@app.route("/")
def index():
    return render_template("index.html")

# 🔥 загрузка pptx
@app.route("/upload", methods=["POST"])
def upload():
    file = request.files["file"]
    file_id = str(uuid.uuid4())
    pptx_path = os.path.join(UPLOAD_FOLDER, file_id + ".pptx")
    file.save(pptx_path)

    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:
                img_path = os.path.join(UPLOAD_FOLDER, f"{file_id}_{i}.png")
                with open(img_path, "wb") as f:
                    f.write(shape.image.blob)
                slides.append(f"{file_id}_{i}.png")

    return jsonify({
        "id": file_id,
        "slides": slides
    })

# 🔥 отдать картинку
@app.route("/image/<name>")
def image(name):
    return send_file(os.path.join(UPLOAD_FOLDER, name))

# 🔥 AI обработка ВСЕХ слайдов
@app.route("/process_ai", methods=["POST"])
def process_ai():
    file_id = request.json["id"]

    processed = []

    for file in os.listdir(UPLOAD_FOLDER):
        if file.startswith(file_id) and file.endswith(".png"):

            path = os.path.join(UPLOAD_FOLDER, file)

            # 👉 отправка в AI
            res = requests.post(
                "https://api.replicate.com/v1/predictions",
                headers={
                    "Authorization": f"Token {API_TOKEN}",
                    "Content-Type": "application/json"
                },
                json={
                    "version": "stability-ai/stable-diffusion-inpainting",
                    "input": {
                        "image": open(path, "rb"),
                        "prompt": "remove watermark, logo, text in bottom right corner, clean background"
                    }
                }
            ).json()

            get_url = res["urls"]["get"]

            while True:
                r = requests.get(get_url, headers={
                    "Authorization": f"Token {API_TOKEN}"
                }).json()

                if r["status"] == "succeeded":
                    output_url = r["output"][0]

                    img_data = requests.get(output_url).content
                    new_path = os.path.join(OUTPUT_FOLDER, file)

                    with open(new_path, "wb") as f:
                        f.write(img_data)

                    processed.append(new_path)
                    break

                if r["status"] == "failed":
                    break

                time.sleep(2)

    # 🔥 сборка обратно в pptx
    prs = Presentation()

    for img in processed:
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(img, 0, 0, width=Inches(10))

    result_path = os.path.join(OUTPUT_FOLDER, file_id + ".pptx")
    prs.save(result_path)

    return jsonify({"done": True})

# 🔥 скачать
@app.route("/download/<file_id>")
def download(file_id):
    return send_file(
        os.path.join(OUTPUT_FOLDER, file_id + ".pptx"),
        as_attachment=True
    )

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
